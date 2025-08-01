import os
import fitz  # PyMuPDF for PDFs
import docx  # python-docx for Word files
import pptx  # python-pptx for PowerPoint
import ollama
import json
from flask import Flask, request, jsonify

app = Flask(__name__)

DOCUMENTS_DIR = "documents"
os.makedirs(DOCUMENTS_DIR, exist_ok=True)  # Ensure the directory exists

# Function to extract text from PDFs
def extract_text_from_pdf(file_path):
    text = ""
    try:
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text("text") + "\n"
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
    return text

# Function to extract text from Word documents
def extract_text_from_docx(file_path):
    text = ""
    try:
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
    return text

# Function to extract text from PowerPoint slides
def extract_text_from_pptx(file_path):
    text = ""
    try:
        presentation = pptx.Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
    return text

# Function to extract text based on file type
def extract_text(file_path):
    if file_path.endswith(".pdf"):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith(".docx"):
        return extract_text_from_docx(file_path)
    elif file_path.endswith(".pptx"):
        return extract_text_from_pptx(file_path)
    else:
        return ""  # Skip non-readable files

# Function to load all documents dynamically
def load_documents():
    documents = []
    for file in os.listdir(DOCUMENTS_DIR):
        file_path = os.path.join(DOCUMENTS_DIR, file)
        text = extract_text(file_path)
        if text.strip():  # Only add files with actual content
            documents.append({"file": file, "content": text})
    return documents

# Function to rank results based on query relevance
def rank_results(results, query):
    ranked = []
    for result in results:
        text = result.get("content", "")
        file_name = result.get("file", "Unknown File")

        prompt = f"Rate how relevant the following text is to this query: {query}. Provide ONLY a score from 0 to 100.\n\nText: {text[:1000]}"  # Limit to first 1000 chars
        response = ollama.chat(model="gemma2:2b", messages=[{"role": "user", "content": prompt}])

        try:
            score = int(response.get("message", {}).get("content", "0").strip())
        except ValueError:
            score = 0

        ranked.append({"score": score, "file": file_name, "content": text})

    ranked = [doc for doc in ranked if doc["score"] > 0]  # Remove zero-score docs
    ranked.sort(reverse=True, key=lambda x: x["score"])
    return ranked

# Function to generate summaries
def generate_summary(text, query):
    prompt = f"Summarize the following document based on how it relates to the query: {query}\n\nText: {text[:2000]}"  # Limit to first 2000 chars
    response = ollama.chat(model="gemma2:2b", messages=[{"role": "user", "content": prompt}])
    return response.get("message", {}).get("content", "Summary not available.")

@app.route("/search", methods=["POST"])
def search_documents():
    data = request.json
    query = data.get("query", "").strip()

    if not query:
        return jsonify({"error": "Query cannot be empty"}), 400

    documents = load_documents()
    ranked_results = rank_results(documents, query)

    for result in ranked_results:
        result["summary"] = generate_summary(result["content"], query)

    return jsonify(ranked_results[:5])

@app.route("/upload", methods=["POST"])
def upload_file():
    if "file" not in request.files:
        return jsonify({"error": "No file part"}), 400

    file = request.files["file"]
    
    if file.filename == "":
        return jsonify({"error": "No selected file"}), 400

    file_path = os.path.join(DOCUMENTS_DIR, file.filename)
    file.save(file_path)
    
    return jsonify({"message": "File uploaded successfully", "file": file.filename})

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)

